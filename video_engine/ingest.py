#!/usr/bin/env python3
"""教材檔案 → 純文字。

引擎的輸入契約其實是「一段純文字」——generate_lesson.py 只做 open().read()，
沒有任何 markdown 解析。所以支援新格式就是在最前面多一道抽取，
後面七個階段一行都不用改。
"""
import hashlib
import os
import re

SUPPORTED = (".md", ".txt", ".pptx")


def lesson_id_for(path):
	"""教材路徑 → lesson_id。run.py 與服務層共用這一份，各算一次遲早走鐘"""
	stem = os.path.splitext(os.path.basename(path))[0].lower()
	slug = re.sub(r"[^a-z0-9_]", "_", stem).strip("_")
	# 中文檔名會被整串壓成底線，不同檔案全部撞在同一個 id 上。
	# 退化時改用原檔名的雜湊，至少能區分
	if not slug.replace("_", ""):
		return "material_" + hashlib.sha1(stem.encode("utf-8")).hexdigest()[:8]
	return slug


def _pptx_text(path):
	"""每頁抓標題、內文與備忘稿。

	投影片順序、備忘稿歸屬都交給 python-pptx——這兩件事在檔案裡都隔了
	一層 rels 間接，照檔名編號硬對會張冠李戴。備忘稿常常是老師真正想講的話，
	比投影片上的關鍵字有用。
	"""
	from pptx import Presentation

	blocks = []
	for num, slide in enumerate(Presentation(path).slides, start=1):
		title = slide.shapes.title
		# 用 shape_id 比對，不能用 is：slide.shapes 每次迭代都給新的 proxy 物件
		title_id = title.shape_id if title is not None else None
		head = title.text.strip() if title is not None and title.has_text_frame else ""
		lines = []
		for shape in slide.shapes:
			if shape.shape_id == title_id or not shape.has_text_frame:
				continue
			for para in shape.text_frame.paragraphs:
				text = "".join(r.text for r in para.runs).strip()
				if text:
					lines.append(text)
		if not (head or lines):
			continue
		block = [f"## 第 {num} 頁：{head}" if head else f"## 第 {num} 頁"]
		block += [f"- {line}" for line in lines]
		if slide.has_notes_slide:
			said = slide.notes_slide.notes_text_frame.text.strip()
			if said:
				block.append("備忘稿：" + " ".join(said.split()))
		blocks.append("\n".join(block))
	if not blocks:
		raise ValueError("這份簡報沒有文字層，抽不出任何內容")
	return "\n\n".join(blocks)


def extract_text(path):
	"""教材檔案 → 純文字。副檔名不支援或抽不到內容時丟 ValueError"""
	ext = os.path.splitext(path)[1].lower()
	if ext not in SUPPORTED:
		raise ValueError(f"不支援的格式 {ext}，只吃 {'、'.join(SUPPORTED)}")
	if ext == ".pptx":
		return _pptx_text(path)
	return open(path, encoding="utf-8").read()
