#!/usr/bin/env python3
"""產測試用的簡報。

用 python-pptx 產，不手工組 XML：手工檔案沒有 presentation.xml 的播放順序表
與 _rels 關聯鏈，測不出「備忘稿檔名編號 != 投影片編號」這類真實陷阱。
"""
from pptx import Presentation


def make(path, slides, notes=None):
	"""slides 是 [(標題, [內文, ...]), ...]；notes 是 {投影片編號(1 起): 備忘稿}"""
	notes = notes or {}
	prs = Presentation()
	for i, (title, bullets) in enumerate(slides, start=1):
		slide = prs.slides.add_slide(prs.slide_layouts[1])
		slide.shapes.title.text = title
		slide.placeholders[1].text = "\n".join(bullets)
		if i in notes:
			slide.notes_slide.notes_text_frame.text = notes[i]
	prs.save(path)


def make_blank(path, pages=1):
	"""沒有任何文字的簡報，用來測「全是圖」的失敗路徑"""
	prs = Presentation()
	for _ in range(pages):
		prs.slides.add_slide(prs.slide_layouts[6])   # 空白版面
	prs.save(path)
